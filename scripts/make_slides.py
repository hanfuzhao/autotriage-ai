"""Builds the pitch deck, PITCH.pptx, from a fixed set of slides."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

BASE = Path(__file__).resolve().parent.parent
OUT = BASE / "PITCH.pptx"

BG = RGBColor(0x0B, 0x10, 0x20)
CARD = RGBColor(0x16, 0x1F, 0x38)
TEXT = RGBColor(0xE8, 0xED, 0xF7)
MUTED = RGBColor(0x93, 0xA0, 0xBD)
ACCENT = RGBColor(0x4C, 0xC9, 0xF0)
ACCENT2 = RGBColor(0x43, 0x61, 0xEE)
GOOD = RGBColor(0x2D, 0xD4, 0xA7)

W, H = Inches(13.333), Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def add_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def line(tf, text, size, color=TEXT, bold=False, first=False, align=PP_ALIGN.LEFT, space=10):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return p


def accent_bar(slide, top=Inches(0.7)):
    bar = slide.shapes.add_shape(1, Inches(0.9), top, Inches(0.16), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return bar


def title_slide(prs):
    s = add_slide(prs)
    accent = s.shapes.add_shape(1, 0, Inches(2.5), Inches(0.28), Inches(2.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT; accent.line.fill.background()
    tf = box(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(3.0))
    line(tf, "AutoTriage AI", 54, TEXT, bold=True, first=True, space=6)
    line(tf, "Is your car problem a safety issue?", 26, ACCENT, space=24)
    line(tf, "Module 2 Project, Natural Language Processing", 18, MUTED, space=4)
    line(tf, "Live demo: autotriage-ai-unfvnsiy6a-uc.a.run.app", 16, MUTED)
    return s


def content_slide(prs, title, bullets, sub=None):
    s = add_slide(prs)
    accent_bar(s)
    tf = box(s, Inches(1.2), Inches(0.62), Inches(11), Inches(1.0))
    line(tf, title, 34, TEXT, bold=True, first=True)
    if sub:
        line(tf, sub, 18, MUTED, space=4)
    body = box(s, Inches(1.2), Inches(1.95), Inches(11), Inches(5.0))
    for i, b in enumerate(bullets):
        p = line(body, b, 22, TEXT, first=(i == 0), space=18)
        p.level = 0
    return s


def results_slide(prs):
    s = add_slide(prs)
    accent_bar(s)
    tf = box(s, Inches(1.2), Inches(0.62), Inches(11), Inches(1.0))
    line(tf, "Results on the held-out test set", 34, TEXT, bold=True, first=True)

    rows = [
        ("Model", "Accuracy", "Macro-F1"),
        ("Naive, majority class", "0.094", "0.012"),
        ("Classical, TF-IDF and LogReg", "0.774", "0.768"),
        ("Deep, TextCNN and GloVe, deployed", "0.771", "0.762"),
    ]
    table = s.shapes.add_table(len(rows), 3, Inches(1.6), Inches(2.2), Inches(10), Inches(3.0)).table
    table.columns[0].width = Inches(5.6)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.2)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT2 if r == 0 else CARD
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = para.add_run(); run.text = val
            run.font.size = Pt(18 if r else 18)
            run.font.bold = (r == 0) or (r == 3)
            run.font.color.rgb = TEXT if r else RGBColor(0xFF, 0xFF, 0xFF)
    note = box(s, Inches(1.6), Inches(5.6), Inches(10), Inches(1.4))
    line(note, "Long-tailed data, so macro-F1 is the metric that matters. The naive floor is 0.012.",
         18, MUTED, first=True, space=6)
    line(note, "GloVe transfer learning pulls the deep model up to the classical baseline.", 18, MUTED)
    return s


def build():
    prs = new_deck()
    title_slide(prs)
    content_slide(prs, "The problem", [
        "Your car does something odd, a noise, a warning light, a scary moment",
        "Most owners cannot tell a real safety defect from a harmless quirk",
        "So they worry over nothing, or keep driving through a real problem",
        "Recalls start from owner complaints, but only if people recognise and report them",
    ])
    content_slide(prs, "What I built", [
        "Describe the problem in plain words, the way you would tell a friend",
        "It tells you which of 14 car systems it is, and how serious that system is",
        "And whether it is worth reporting to NHTSA",
        "It shows the exact words behind the call, so you can decide to trust it",
    ])
    content_slide(prs, "Why this is new work", [
        "Not my hackathon: that was sentiment on consumer reviews, fine-tuned DistilBERT",
        "This is safety complaints, which system failed, a TextCNN with GloVe",
        "Published work on this database clusters it for regulators (Ghazizadeh et al., 2014)",
        "I pointed the same data at the owner, with explanations and a safety tier",
    ], sub="Different data, question, model, and user")
    content_slide(prs, "Live demo", [
        "autotriage-ai-unfvnsiy6a-uc.a.run.app",
        "Most likely system, confidence, and the other systems it could be",
        "The words that drove the call, plus a safety tier and a plain next step",
        "Flip between the naive, classical, and deep models on the same text",
    ], sub="Let the screen do the talking")
    content_slide(prs, "Three models, judged honestly", [
        "Naive baseline, majority class, the accuracy floor",
        "Classical, TF-IDF with class-weighted Logistic Regression",
        "Deep, TextCNN with GloVe transfer learning, the deployed model",
        "The data is long-tailed, so we judge on macro-F1, not accuracy",
    ])
    results_slide(prs)
    content_slide(prs, "What I learned", [
        "GloVe transfer learning erased the deep model's cold-start gap",
        "At about 1,300 examples it went from 11 points behind to a tie",
        "For a safety tool, honesty matters: it says 'not sure' instead of a false all-clear",
        "Answering only the most confident 60 percent reaches 91 percent accuracy",
    ])
    content_slide(prs, "Why it matters, and the ask", [
        "Turn a helpless worry into a clear read and a next step",
        "Every good report an owner files also feeds the recall system",
        "Live on the cloud, a small model, no GPU, and it explains itself",
        "The ask: put it in a car marketplace or claims flow and test it with real owners",
    ])
    prs.save(OUT)
    print(f"wrote {OUT} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    build()
